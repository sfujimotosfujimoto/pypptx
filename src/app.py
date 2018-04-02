import os
import sys
from pandas import read_csv
import tkinter as tk
from tkinter import filedialog as fd
from pptx import Presentation


ppt_base = ''

if getattr(sys, 'frozen', False):
    bundle_dir = sys._MEIPASS
    ppt_base = os.path.abspath(bundle_dir + '/base.pptx')
else:
    bundle_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    ppt_base = os.path.abspath(bundle_dir + '/ppt/base.pptx')


print(f"sysargv[0]: {sys.argv[0]}")
print(f"bundle_dir: {bundle_dir}")
print(f"ppt_base: {ppt_base}")

path = os.path.dirname(os.path.abspath(__file__))
print(f"currentdir: {path}")
print("Where is your data?")
print("You can drag your file to this terminal.")
# from_path = input().strip()
# print(f"from_path: {from_path}")
print("Where should I put your data?")
print("Don't forget to add `.pptx` to the name of the file")
# to_path = input().strip()
# print(f"to_path: {to_path}")

prs = Presentation(os.path.abspath(ppt_base))

filename = ""


def get_path(event):
    global filename
    filename = fd.askopenfilename(initialdir=".", title="select file")
    to_entry_input = to_entry.get()
    print(to_entry_input)
    print(filename)


def create_pptx(event):
    global filename
    to_entry_input = to_entry.get()
    dirname = os.path.dirname(os.path.abspath(filename))
    to_name = f"{dirname}/{to_entry_input}.pptx"
    count = 0
    print("Created")
    df = read_csv(os.path.abspath(filename))
    for index, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        try:
            slide.placeholders[0].text = row[0]
            slide.placeholders[1].text = row[1]
            if index % 10 == 0:
                print(f"Made {index+1} slides.")
            count = index
        except ValueError:
            pass
        except TypeError:
            pass
    prs.save(os.path.abspath(to_name))
    print(
        f"Saved! Made a total of {count} slides!\nYour files is at: {to_name}")
    to_entry.delete(0, 'end')


root = tk.Tk()
root.geometry("300x200")

from_label = tk.Label(root, text="Where is your data?")
from_button = tk.Button(root, text="Enter")
from_button.bind("<Button-1>", get_path)

to_label = tk.Label(root, text="What should I name your file?")
to_entry = tk.Entry(root)
to_entry_input = to_entry.get()
print(to_entry_input)
to_button = tk.Label(root, text="Create")
to_button.bind("<Button-1>", create_pptx)

from_label.pack()
from_button.pack()

to_label.pack()
to_entry.pack()
to_button.pack()

root.mainloop()
